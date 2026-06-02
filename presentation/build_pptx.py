#!/usr/bin/env python3
"""
Build the intermediate-status presentation (max 4 slides) as a .pptx.

Run:  python3 presentation/build_pptx.py
Out:  presentation/Bierkastenerkennung_Zwischenstand.pptx

Content reflects the real repo state (Kastendetektion + Warp/Grid done,
Stage 1/2 + training pending) and the professor's brief: what/how/why,
difficulties, June plan. Demo images come from demo_output/.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "demo_output"
OUT = ROOT / "presentation" / "Bierkastenerkennung_Zwischenstand.pptx"

# Palette
NAVY = RGBColor(0x1B, 0x3A, 0x5B)
NAVY_DARK = RGBColor(0x14, 0x2C, 0x46)
AMBER = RGBColor(0xE1, 0xA1, 0x00)
GREY_BG = RGBColor(0xF4, 0xF6, 0xF8)
GREY_BOX = RGBColor(0xD7, 0xDD, 0xE3)
INK = RGBColor(0x21, 0x25, 0x29)
MUTE = RGBColor(0x6B, 0x74, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x9E, 0x4F)
RED = RGBColor(0xC0, 0x39, 0x2B)

SW, SH = Inches(13.333), Inches(7.5)
TEAM = "Mia Scharpf · David Krieger · Robin Müller — AKI"


def _no_line(shape) -> None:
    shape.line.fill.background()


def fill_rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    _no_line(sp)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    return tb, tf


def set_run(p, text, *, size, bold=False, color=INK, italic=False, font="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def bullets(tf, items, *, size=15, color=INK, space_after=8, bullet_color=AMBER):
    """items: list of (text, bold) or (text, bold, sub_bool)."""
    first = True
    for item in items:
        text, bold = item[0], item[1]
        sub = item[2] if len(item) > 2 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        glyph = "–  " if sub else "•  "
        gr = p.add_run()
        gr.text = glyph
        gr.font.size = Pt(size)
        gr.font.bold = True
        gr.font.color.rgb = MUTE if sub else bullet_color
        gr.font.name = "Calibri"
        set_run(p, text, size=size, bold=bold, color=(MUTE if sub else color))


def header(slide, title, step_label):
    fill_rect(slide, 0, 0, SW, Inches(1.02), NAVY)
    fill_rect(slide, 0, Inches(1.02), SW, Pt(4), AMBER)
    tb, tf = textbox(slide, Inches(0.55), Inches(0.12), Inches(9.6), Inches(0.8),
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_run(p, title, size=27, bold=True, color=WHITE)
    # step badge
    badge = fill_rect(slide, Inches(10.5), Inches(0.27), Inches(2.4), Inches(0.5),
                      AMBER, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf2 = badge.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2, step_label, size=12, bold=True, color=NAVY_DARK)


def footer(slide, page):
    fill_rect(slide, Inches(0.55), Inches(7.02), Inches(12.23), Pt(1), GREY_BOX)
    tb, tf = textbox(slide, Inches(0.55), Inches(7.08), Inches(9), Inches(0.34))
    p = tf.paragraphs[0]
    set_run(p, TEAM, size=9, color=MUTE)
    tb2, tf2 = textbox(slide, Inches(10.3), Inches(7.08), Inches(2.45), Inches(0.34))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    set_run(p2, f"Zwischenstand · 02.06.2026 · {page}/4", size=9, color=MUTE)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_rect(s, 0, 0, SW, SH, WHITE)  # solid white background
    return s


def add_caption(slide, x, y, w, text):
    tb, tf = textbox(slide, x, y, w, Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, text, size=10, italic=True, color=MUTE)


# ----------------------------------------------------------------------------
def slide1_concept(prs):
    s = blank(prs)
    # title band
    fill_rect(s, 0, 0, SW, Inches(2.05), NAVY)
    fill_rect(s, 0, Inches(2.05), SW, Pt(5), AMBER)
    tb, tf = textbox(s, Inches(0.6), Inches(0.42), Inches(12), Inches(1.5))
    p = tf.paragraphs[0]
    set_run(p, "Bierkasten-Erkennung", size=40, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    set_run(p2, "Zwischenstand — Hybridansatz: Klassische Bildverarbeitung + Machine Learning",
            size=18, color=RGBColor(0xCF, 0xDD, 0xEC))
    p3 = tf.add_paragraph()
    p3.space_before = Pt(8)
    set_run(p3, TEAM, size=13, color=AMBER, bold=True)

    # left: goal + questions
    tb, tf = textbox(s, Inches(0.6), Inches(2.45), Inches(6.0), Inches(3.0))
    p = tf.paragraphs[0]
    set_run(p, "Ziel", size=16, bold=True, color=NAVY)
    p = tf.add_paragraph()
    p.space_after = Pt(10)
    set_run(p, "Füllzustand eines Bierkastens in Echtzeit per Kamera erkennen "
               "(Prototyp: Paulaner Weißbier, 4×5 = 20 Flaschen).", size=14, color=INK)
    bullets(tf, [
        ("Ist der Kasten vollständig befüllt? (Flasche vorhanden / Slot leer)", False),
        ("Sind die Flaschen voll oder leer? (Kronkorken vorhanden = voll)", False),
    ], size=14, space_after=7)

    # right: why hybrid
    box = fill_rect(s, Inches(6.95), Inches(2.45), Inches(5.78), Inches(2.05),
                    GREY_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    box.line.color.rgb = GREY_BOX
    box.line.width = Pt(1)
    tb, tf = textbox(s, Inches(7.2), Inches(2.62), Inches(5.3), Inches(1.8))
    p = tf.paragraphs[0]
    set_run(p, "Warum hybrid?", size=16, bold=True, color=NAVY)
    bullets(tf, [
        ("Klassische CV: schnell, deterministisch, kein Training", False),
        ("ML: robust gegen Licht, Reflexionen, Perspektive", False),
        ("→ Geometrie klassisch, Semantik per ML", True),
    ], size=13, space_after=6)

    # pipeline chevrons
    tb, tf = textbox(s, Inches(0.6), Inches(4.95), Inches(8), Inches(0.35))
    set_run(tf.paragraphs[0], "Pipeline (pro Frame)", size=14, bold=True, color=NAVY)

    steps = ["Aufnahme", "Kasten-\ndetektion", "Entzerrung", "Grid 4×5",
             "Stufe 1\nvoll/leer", "Stufe 2\nKronkorken", "Overlay +\nStatistik"]
    done = [True, True, True, True, False, False, False]
    x0 = Inches(0.55)
    cw = Inches(1.78)
    gap = Inches(-0.06)
    y = Inches(5.4)
    for i, (label, ok) in enumerate(zip(steps, done)):
        x = Emu(int(x0) + i * (int(cw) + int(gap)))
        ch = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, cw, Inches(0.92))
        ch.fill.solid()
        ch.fill.fore_color.rgb = AMBER if ok else GREY_BOX
        _no_line(ch)
        ch.shadow.inherit = False
        tf = ch.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for j, line in enumerate(label.split("\n")):
            pp = p if j == 0 else tf.add_paragraph()
            pp.alignment = PP_ALIGN.CENTER
            set_run(pp, line, size=11, bold=True, color=(NAVY_DARK if ok else MUTE))

    # legend for chevrons
    lx = Inches(0.6)
    ly = Inches(6.5)
    fill_rect(s, lx, ly, Inches(0.28), Inches(0.18), AMBER)
    tb, tf = textbox(s, Emu(int(lx) + int(Inches(0.36))), Emu(int(ly) - int(Inches(0.04))),
                     Inches(2.0), Inches(0.3))
    set_run(tf.paragraphs[0], "umgesetzt", size=11, color=INK)
    lx2 = Inches(2.4)
    fill_rect(s, lx2, ly, Inches(0.28), Inches(0.18), GREY_BOX)
    tb, tf = textbox(s, Emu(int(lx2) + int(Inches(0.36))), Emu(int(ly) - int(Inches(0.04))),
                     Inches(2.0), Inches(0.3))
    set_run(tf.paragraphs[0], "geplant (Juni)", size=11, color=INK)

    notes(s, "Begrüßung und Projektziel: Wir erkennen in Echtzeit, wie voll ein Bierkasten ist – "
             "zwei Fragen: Sind alle Flaschen da, und sind sie voll (Kronkorken) oder leer. "
             "Prototyp ist der Paulaner-Kasten mit 20 Slots im 4x5-Raster. "
             "Methodisch kombinieren wir klassische Bildverarbeitung (schnell, kein Training) mit "
             "Machine Learning (robust gegen Licht und Reflexionen). Die Pipeline unten zeigt die "
             "sieben Schritte pro Frame – orange ist umgesetzt, grau ist für Juni geplant. "
             "Wir stehen bei etwa der Hälfte: die ersten vier geometrischen Schritte laufen.")
    footer(s, 1)
    return s


def slide2_done(prs):
    s = blank(prs)
    header(s, "Umgesetzt: Kastendetektion & Entzerrung", "Schritte 2–4")

    tb, tf = textbox(s, Inches(0.55), Inches(1.35), Inches(6.6), Inches(5.4))
    p = tf.paragraphs[0]
    set_run(p, "Was läuft (getestet)", size=15, bold=True, color=NAVY)
    p.space_after = Pt(6)
    bullets(tf, [
        ("Kastendetektion zweigleisig", True),
        ("YOLOv8 (Ultralytics) als primärer Detektor – Fine-Tuning vorbereitet", True),
        ("Klassischer Fallback: Canny + Konturen + minAreaRect", True),
        ("Schnittstelle detect_crate(frame) → Bounding Box + 4 Eckpunkte + Konfidenz", False),
        ("Stärkerer klassischer Detektor: HSV-Rot + Otsu + Hough + Hu-Momente", False),
        ("Perspektivische Entzerrung (Warp) → normierte Draufsicht", False),
        ("Grid Mapping 4×5 → 20 Slot-Mittelpunkte + Rückprojektion + ROI je Slot", False),
        ("Komplette Daten-/Trainings-Pipeline: Label-UI, Split, YOLOv8-Training", False),
        ("End-to-End-Smoke-Test grün: Detektion → Entzerrung → Grid", True),
    ], size=13.5, space_after=6)

    # images side by side (top of right column)
    iw = Inches(2.85)
    x1 = Inches(7.15)
    x2 = Inches(10.05)
    iy = Inches(1.45)
    s.shapes.add_picture(str(IMG / "01_detection.png"), x1, iy, width=iw)
    s.shapes.add_picture(str(IMG / "04_grid_on_original.png"), x2, iy, width=iw)
    add_caption(s, x1, Inches(3.62), iw, "Detektion: Box + 4 Ecken")
    add_caption(s, x2, Inches(3.62), iw, "Warp + 4×5-Grid (20 Slots)")
    # honesty note centered across both
    tb, tf = textbox(s, Inches(7.15), Inches(3.98), Inches(5.75), Inches(0.35))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "Bilder: synthetischer Test — echte Fotos folgen im Juni", size=9.5,
            italic=True, color=MUTE)
    # interface box (fills the lower right, shows the handover contract)
    box = fill_rect(s, Inches(7.15), Inches(4.5), Inches(5.75), Inches(1.95),
                    GREY_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    box.line.color.rgb = GREY_BOX
    box.line.width = Pt(1)
    fill_rect(s, Inches(7.15), Inches(4.5), Pt(6), Inches(1.95), AMBER)
    tb, tf = textbox(s, Inches(7.45), Inches(4.68), Inches(5.3), Inches(1.6))
    p = tf.paragraphs[0]
    set_run(p, "Schnittstelle ans Grid-Modul", size=14, bold=True, color=NAVY)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    set_run(p2, "detect_crate(frame)", size=13, bold=True, color=INK, font="Consolas")
    p3 = tf.add_paragraph()
    p3.space_before = Pt(2)
    set_run(p3, "→ Bounding Box (x, y, w, h)", size=12.5, color=INK, font="Consolas")
    p4 = tf.add_paragraph()
    p4.space_before = Pt(2)
    set_run(p4, "→ 4 Eckpunkte (TL,TR,BR,BL) + Konfidenz", size=12.5, color=INK, font="Consolas")

    notes(s, "Das ist unser aktueller Fokus und Kern dieses Termins. Die Kastendetektion läuft "
             "zweigleisig: YOLOv8 ist der primäre Detektor – das Fine-Tuning ist vorbereitet, aber "
             "noch nicht trainiert; deshalb nutzen wir aktuell den klassischen Fallback (Canny, "
             "Konturen, minAreaRect), der sofort ohne Training funktioniert. Die Schnittstelle "
             "detect_crate liefert eine Bounding Box plus vier Eckpunkte – genau die Eingabe fürs "
             "Grid-Modul. Darauf folgt die perspektivische Entzerrung in eine Draufsicht und das "
             "4x5-Grid mit 20 Slot-Mittelpunkten, die wir auch ins Originalbild zurückrechnen. "
             "Ein automatischer Smoke-Test bestätigt, dass die Kette Detektion → Entzerrung → Grid "
             "durchläuft. Wichtig zur Ehrlichkeit: Die gezeigten Bilder sind ein synthetischer Test, "
             "da wir noch keine echten gelabelten Fotos haben.")
    footer(s, 2)
    return s


def slide3_problems(prs):
    s = blank(prs)
    header(s, "Schwierigkeiten & Erkenntnisse", "Offene Punkte")

    cards = [
        ("Noch keine Trainingsdaten",
         "YOLO läuft bisher nur über den klassischen Fallback. Eigene Fotos, "
         "best.pt und das Ziel mAP > 0,8 stehen noch aus."),
        ("Eckpunkte nur Näherung",
         "YOLO liefert eine achsparallele Box → die 4 Ecken für die Entzerrung sind "
         "approximiert. Orientierung gibt bisher nur der klassische Pfad (OBB/Pose offen)."),
        ("Hauptrisiko: Licht & Reflexionen",
         "Wechselndes Licht und Reflexionen auf Glas/Kronkorken erschweren die spätere "
         "Klassifikation – laut Risikoanalyse das größte Risiko."),
        ("Stufe 1 & 2 noch offen",
         "Slot voll/leer und Kronkorken-Erkennung sind konzipiert, aber noch nicht "
         "implementiert – das ist der Schwerpunkt für Juni."),
    ]
    x = [Inches(0.55), Inches(6.85)]
    y = [Inches(1.45), Inches(4.2)]
    cw, chh = Inches(5.95), Inches(2.5)
    for i, (title, body) in enumerate(cards):
        cx = x[i % 2]
        cy = y[i // 2]
        card = fill_rect(s, cx, cy, cw, chh, GREY_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
        card.line.color.rgb = GREY_BOX
        card.line.width = Pt(1)
        fill_rect(s, cx, cy, Pt(6), chh, AMBER)
        tb, tf = textbox(s, Emu(int(cx) + int(Inches(0.3))), Emu(int(cy) + int(Inches(0.22))),
                         Inches(5.4), Emu(int(chh) - int(Inches(0.4))))
        p = tf.paragraphs[0]
        set_run(p, title, size=16, bold=True, color=NAVY)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        set_run(p2, body, size=13, color=INK)

    notes(s, "Wo hakt es? Erstens: Wir haben noch keine eigenen gelabelten Trainingsbilder, "
             "deshalb läuft die Detektion aktuell über den klassischen Fallback statt über das "
             "trainierte YOLO – das Modell best.pt und das mAP-Ziel über 0,8 stehen noch aus. "
             "Zweitens: YOLO gibt nur eine achsparallele Box aus, daher sind die vier Eckpunkte "
             "für die Entzerrung nur eine Näherung; eine echte gedrehte Box käme bisher nur aus "
             "dem klassischen Pfad. Drittens, und das ist das Hauptrisiko laut unserer Analyse: "
             "wechselndes Licht und Reflexionen auf Glas und Kronkorken. Und viertens sind die "
             "beiden Klassifikationsstufen zwar durchdacht, aber noch nicht umgesetzt – genau "
             "das nehmen wir uns für den Juni vor.")
    footer(s, 3)
    return s


def slide4_outlook(prs):
    s = blank(prs)
    header(s, "Ausblick: bis zur Abschlusspräsentation", "23.06.2026")

    tb, tf = textbox(s, Inches(0.55), Inches(1.35), Inches(7.4), Inches(5.4))
    p = tf.paragraphs[0]
    set_run(p, "Meilensteine Juni", size=15, bold=True, color=NAVY)
    p.space_after = Pt(8)
    bullets(tf, [
        ("Fotos aufnehmen (Winkel/Licht/Rotation) + labeln → YOLOv8 fine-tunen (best.pt, mAP > 0,8)", False),
        ("Stufe 1 – Slot voll/leer: klassisch (Hough/Helligkeit) + ML (CNN / HOG+SVM)", False),
        ("Stufe 2 – Kronkorken: HSV-Fallback + CNN (robust gegen Reflexionen)", False),
        ("Farbcodiertes Overlay + Gesamtstatistik (z. B. „15/20, davon 12 voll“)", False),
        ("Frame-Stabilisierung (10 stabile Frames) gegen Flackern", False),
        ("Genauere Entzerrung über echte Oberflächen-Ecken / Oriented Bounding Box", False),
    ], size=14, space_after=11)

    # legend panel for final overlay colors
    box = fill_rect(s, Inches(8.3), Inches(1.45), Inches(4.45), Inches(2.45),
                    GREY_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    box.line.color.rgb = GREY_BOX
    box.line.width = Pt(1)
    tb, tf = textbox(s, Inches(8.6), Inches(1.62), Inches(3.9), Inches(0.4))
    set_run(tf.paragraphs[0], "Ziel-Overlay (Legende)", size=14, bold=True, color=NAVY)
    legend = [(GREEN, "grün = Flasche voll"), (AMBER, "gelb = Flasche leer"),
              (RED, "rot = Slot fehlt")]
    ly = Inches(2.25)
    for col, label in legend:
        fill_rect(s, Inches(8.65), ly, Inches(0.34), Inches(0.34), col,
                  MSO_SHAPE.OVAL)
        tb, tf = textbox(s, Inches(9.15), Emu(int(ly) - int(Inches(0.02))),
                         Inches(3.4), Inches(0.4))
        set_run(tf.paragraphs[0], label, size=13, color=INK)
        ly = Emu(int(ly) + int(Inches(0.55)))

    # closing strip
    strip = fill_rect(s, Inches(8.3), Inches(4.2), Inches(4.45), Inches(1.7),
                      NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
    tb, tf = textbox(s, Inches(8.6), Inches(4.42), Inches(3.9), Inches(1.3),
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_run(p, "Fundament steht", size=15, bold=True, color=AMBER)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    set_run(p2, "Geometrie-Pipeline (Detektion → Entzerrung → Grid) ist die Basis "
                "für die ML-Klassifikation im Juni.", size=12.5, color=WHITE)

    notes(s, "Unser Plan bis zur Abschlusspräsentation am 23. Juni: Zuerst echte Fotos in "
             "verschiedenen Winkeln, Lichtsituationen und Rotationen aufnehmen und labeln, dann "
             "YOLOv8 fine-tunen mit dem Ziel mAP über 0,8. Darauf bauen die beiden "
             "Klassifikationsstufen auf: Stufe 1 entscheidet je Slot voll oder leer – klassisch "
             "über Hough-Kreise und Helligkeit, plus ein kleines CNN oder HOG+SVM. Stufe 2 prüft "
             "den Kronkorken über HSV als Fallback und ein CNN. Am Ende steht das farbcodierte "
             "Overlay – grün voll, gelb leer, rot fehlt – mit einer Gesamtstatistik, dazu eine "
             "Frame-Stabilisierung gegen Flackern und eine genauere Entzerrung über echte "
             "Oberflächen-Ecken. Wichtig: Die geometrische Pipeline steht bereits – sie ist das "
             "Fundament, auf dem die ML-Klassifikation im Juni aufsetzt. Vielen Dank.")
    footer(s, 4)
    return s


def main() -> None:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    slide1_concept(prs)
    slide2_done(prs)
    slide3_problems(prs)
    slide4_outlook(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Gespeichert: {OUT}")


if __name__ == "__main__":
    main()
